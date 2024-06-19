import json
import pandas as pd
import os
import math
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import zscore, pearsonr, ttest_ind
import math
from math import ceil
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Function to read JSON files and extract metadata
def read_json_files(files):
    columns = ["knee_side",'filename', 'size', 'regions_shape_attributes_name', 'x1', 'y1', 'x2', 'y2', 'x3', 'y3']
    clean_data_list = []
    fehler_id_rechte = [1,748,1174,2007,2008,2009,2010,3048, 74, 1857]
    fehler_id_linke = [270,326,349,389,423, 485,562,634,656,682,706,710,739, 743,762,1106,1115,1559,1785,1855,1856,1834,366,2027, 2820,3016,3727]
    jj = 0
    count = 0
    for file in files:
        with open(file) as f:
            data = json.load(f)
            # Extract metadata
            for i in data['_via_img_metadata'].keys():
                count = count + 1
                if file == 'Student2_rechts.json' and jj == 0:
                    count = 1
                    jj = 1

                filename = data['_via_img_metadata'][i]['filename']
                size = data['_via_img_metadata'][i]['size']
                regions = data['_via_img_metadata'][i]['regions']
                if count in fehler_id_rechte:
                    if file == 'Student1_rechts.json' or file == 'Student2_rechts.json':
                        print(file, filename, count, "The measurement was wrong based on post_it pics")
                        continue

                if count in fehler_id_linke:
                    if file == 'Student1_links.json' or file == 'Student2_links.json':
                        print(file, filename, count, "The measurement was wrong based on post_it pics")
                        continue
                if regions == []:
                    continue
                if data['_via_img_metadata'][i]['regions'][0] is None:
                    #print(file,filename, "region is none", count)
                    regions_shape_attributes_name = data['_via_img_metadata'][i]['regions'][1]['shape_attributes'][
                        'name']
                    x1 = data['_via_img_metadata'][i]['regions'][1]['shape_attributes']['all_points_x'][0]
                    y1 = data['_via_img_metadata'][i]['regions'][1]['shape_attributes']['all_points_y'][0]
                    x2 = data['_via_img_metadata'][i]['regions'][1]['shape_attributes']['all_points_x'][1]
                    y2 = data['_via_img_metadata'][i]['regions'][1]['shape_attributes']['all_points_y'][1]
                    x3 = data['_via_img_metadata'][i]['regions'][1]['shape_attributes']['all_points_x'][2]
                    y3 = data['_via_img_metadata'][i]['regions'][1]['shape_attributes']['all_points_y'][2]
                else:
                    if 'all_points_y' not in data['_via_img_metadata'][i]['regions'][0]['shape_attributes'].keys():
                        print(file,filename, "all_points_y is missing", count)
                        continue
                    regions_shape_attributes_name = data['_via_img_metadata'][i]['regions'][0]['shape_attributes']['name']
                    if len(data['_via_img_metadata'][i]['regions'][0]['shape_attributes']['all_points_x'])<3:
                        print(file,filename, "all_points_x is missing", count)
                        continue
                    x1 = data['_via_img_metadata'][i]['regions'][0]['shape_attributes']['all_points_x'][0]
                    y1 = data['_via_img_metadata'][i]['regions'][0]['shape_attributes']['all_points_y'][0]
                    x2 = data['_via_img_metadata'][i]['regions'][0]['shape_attributes']['all_points_x'][1]
                    y2 = data['_via_img_metadata'][i]['regions'][0]['shape_attributes']['all_points_y'][1]
                    x3 = data['_via_img_metadata'][i]['regions'][0]['shape_attributes']['all_points_x'][2]
                    y3 = data['_via_img_metadata'][i]['regions'][0]['shape_attributes']['all_points_y'][2]

                # Calculate IS ratio
                # Step 1: Calculate Patellar Length (PL)
                PL = math.sqrt((x3 - x2) ** 2 + (y3 - y2) ** 2)
                # Step 2: Calculate Patellar Tendon Length (TL)
                TL = math.sqrt((x2 - x1) ** 2 + (y2 - y1) ** 2)
                # Step 3: Compute Insall-Salvati Index
                Insall_Salvati_Index = TL / PL
                # Remove part before "_" and after "."
                folder_name_cleaned = file.split('_')[-1].split('.')[0]
                # Remove "CR." from the beginning and ".png" from the end
                file_name_cleaned = filename.replace('CR.', '').replace('.png', '')
                clean_data_list.append([folder_name_cleaned, file_name_cleaned, size, regions_shape_attributes_name, x1, y1, x2, y2, x3, y3, Insall_Salvati_Index])


    df = pd.DataFrame(clean_data_list, columns=columns + ['IS_ratio'])
    return df


# Function to analyze the data
def analyze_data(df):
    desc_stats = df.describe()
    numeric_cols = df.select_dtypes(include=[float, int]).columns
    median = df[numeric_cols].median()
    median_df = pd.DataFrame(median).transpose()
    median_df.index = ['median']
    desc_stats = pd.concat([desc_stats, median_df])
    desc_stats.to_csv('descriptive_statistics_before_outliers_removal.csv')
    print("Descriptive statistics saved to descriptive_statistics_before_outliers_removal.csv")

    pr = len(df)
    df['z_score'] = zscore(df['IS_ratio'])
    z_lower = -3
    z_upper = 3
    outliers = df[(df['z_score'] < z_lower) | (df['z_score'] > z_upper)]
    df = df[(df['z_score'] > z_lower) & (df['z_score'] < z_upper)]
    df = df.drop(columns=['z_score'])
    print("Number of removed outliers after Z-score =", pr - len(df))

    desc_stats = df.describe()
    numeric_cols = df.select_dtypes(include=[float, int]).columns
    median = df[numeric_cols].median()
    median_df = pd.DataFrame(median).transpose()
    median_df.index = ['median']
    desc_stats = pd.concat([desc_stats, median_df])
    desc_stats.to_csv('Task_B_descriptive_statistics_after_outliers_removal.csv')
    print("Descriptive statistics saved to Task_B_descriptive_statistics_after_outliers_removal.csv")

    plt.figure(figsize=(10, 6))
    sns.boxplot(x=df['IS_ratio'])
    plt.title('Boxplot of Insall-Salvati Ratio')
    plt.savefig('boxplot_is_ratio.png')
    plt.show()

    plt.figure(figsize=(10, 6))
    sns.histplot(df['IS_ratio'], kde=True, bins=30)
    plt.title('Histogram of Insall-Salvati Ratio')
    plt.savefig('histogram_is_ratio.png')
    plt.show()
    return df, desc_stats, outliers

"""
    coords = ['x1', 'y1', 'x2', 'y2', 'x3', 'y3']
    for coord in coords:
        plt.figure(figsize=(10, 6))
        sns.histplot(df[coord], kde=True, bins=30)
        plt.title(f'Histogram of {coord}')
        plt.savefig(f'histogram_{coord}.png')
        plt.show()
"""

# Function to merge dataframes
def merge_dataframes(df1, df2):
    merged_df = pd.merge(df1, df2, left_on=['filename', 'knee_side'], right_on=['knee_lat', 'op_seite'])
    merged_df = merged_df.drop_duplicates()
    return merged_df


def perform_two_sample_ttest(merged_df, alpha=0.05):
    # Separate the IS Ratios by patientsex
    male_is_ratios = merged_df[merged_df['patientsex'] == 'M']['IS_ratio']
    female_is_ratios = merged_df[merged_df['patientsex'] == 'F']['IS_ratio']

    # Perform two-sample t-test
    t_stat, p_value = ttest_ind(male_is_ratios, female_is_ratios, nan_policy='omit')

    print(f'Two-sample t-test results:')
    print(f'T-statistic: {t_stat:.2f}, P-value: {p_value:.2e}')

    # Decision making
    if p_value <= alpha:
        print(
            f'Reject the null hypothesis (H0). There is a significant difference between the IS Ratios of males and females (p-value ≤ {alpha}).')
        print(f"Mean IS Ratio for males: {male_is_ratios.mean():.2f}")
        print(f"Mean IS Ratio for females: {female_is_ratios.mean():.2f}")
    else:
        print(
            f'Do not reject the null hypothesis (H0). There is no significant difference between the IS Ratios of males and females (p-value > {alpha}).')

    return t_stat, p_value, male_is_ratios.mean(), female_is_ratios.mean()

def explore_relationships(merged_df):
    merged_df = merged_df.drop(columns=['patientsex'])
    if 'age' in merged_df.columns:
        merged_df['age'] = merged_df['age'].apply(ceil).astype(int)

    correlation_results = {}

    for column in merged_df.select_dtypes(include=['object']).columns:
        merged_df[column] = pd.to_numeric(merged_df[column], errors='coerce')

    correlation_matrix = merged_df.corr()
    correlation_matrix.to_excel('correlation_matrix_merged_data.xlsx')
    print("Correlation matrix saved to correlation_matrix_merged_data.xlsx")

    high_correlation_threshold = 0.05
    is_ratio_corr = correlation_matrix['IS_ratio'].abs()
    high_corr_columns = is_ratio_corr[is_ratio_corr > high_correlation_threshold].index

    filtered_corr_matrix = correlation_matrix.loc[high_corr_columns, high_corr_columns]
    filtered_corr_matrix.to_excel('filtered_correlation_matrix_merged_data.xlsx')
    print("Filtered correlation matrix saved to filtered_correlation_matrix_merged_data.xlsx")

    plt.figure(figsize=(12, 10))
    sns.heatmap(filtered_corr_matrix, annot=True, fmt='.2f', cmap='coolwarm', cbar=True)
    plt.title('Filtered Heatmap of Correlation Matrix')
    plt.savefig('filtered_correlation_matrix_heatmap.png')
    plt.show()

    for column in high_corr_columns:
        if column != 'IS_ratio':
            clean_df = merged_df[[column, 'IS_ratio']].dropna()
            if not clean_df.empty:
                plt.figure(figsize=(10, 6))
                sns.scatterplot(x=clean_df[column], y=clean_df['IS_ratio'])
                plt.title(f'Relationship between {column} and IS_ratio')
                plt.savefig(f'relationship_{column}_is_ratio.png')
                plt.show()
                corr, p_value = pearsonr(clean_df[column], clean_df['IS_ratio'])
                correlation_results[column] = {'correlation': corr, 'p_value': p_value}
                print(f'Correlation between {column} and IS_ratio: {corr:.2f}, p-value: {p_value:.2e}')
                if p_value <= 0.05:
                    print(f'The correlation between {column} and IS_ratio is significant.')
                else:
                    print(f'The correlation between {column} and IS_ratio is not significant.')

    return correlation_results






################## Presentation

def create_presentation(desc_stats, outliers, correlation_results, t_test_results):
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Statistical Analysis Results"
    subtitle.text = "Descriptive statistics, correlations, and t-test results"

    # Descriptive Statistics Slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Descriptive Statistics"
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = str(desc_stats)

    # Outliers Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Outliers"
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = str(outliers)

    # Correlation Results Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Correlation Results"
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for column, result in correlation_results.items():
        p = tf.add_paragraph()
        p.text = f"Column: {column}, Correlation: {result['correlation']:.2f}, P-value: {result['p_value']:.2e}"
        if result['p_value'] <= 0.05:
            p.text += " (Significant)"
        else:
            p.text += " (Not significant)"

    # T-Test Results Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Two-Sample T-Test Results"
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    t_stat, p_value, mean_male, mean_female = t_test_results
    tf.text = f"T-statistic: {t_stat:.2f}\nP-value: {p_value:.2e}\n"
    if p_value <= 0.05:
        tf.text += f"Reject the null hypothesis (H0). There is a significant difference between the IS Ratios of males and females (p-value ≤ 0.05).\n"
        tf.text += f"Mean IS Ratio for males: {mean_male:.2f}\nMean IS Ratio for females: {mean_female:.2f}"
    else:
        tf.text += "Do not reject the null hypothesis (H0). There is no significant difference between the IS Ratios of males and females (p-value > 0.05)."

    # Save Presentation
    prs.save('Task_C_statistical_analysis_results.pptx')
    print("Presentation saved as statistical_analysis_results.pptx")

##############

def main():
    # Define the names of the JSON files with their absolute paths
    jfiles = ['Student1_rechts.json', 'Student1_links.json', 'Student2_rechts.json', 'Student2_links.json']
    # Read JSON files and extract metadata
    final_df = read_json_files(jfiles)
    final_df, desc_stats, outliers = analyze_data(final_df)

    # Save to a single Excel file
    output_file_name = 'Task_A_cleaned_combined_measurements_and_IS_ratio.xlsx'
    final_df.to_excel(output_file_name, index=False)

    patient_data_file_name = 'patientdata_anonymized.xlsx'
    patient_df = pd.read_excel(patient_data_file_name)

    merged_df = merge_dataframes(final_df, patient_df)
    merged_df = merged_df.drop(columns=['x1', 'x2', 'x3', 'y1', 'y2', 'y3', 'knee_side', 'op_seite'])
    merged_df = merged_df.drop(columns=['filename', 'bmi', 'pre_LEN_MECH_TIBIA_cm', 'size', 'knee_lat', 'regions_shape_attributes_name'])

    merged_output_file_name = 'merged_data_with_patient_info.xlsx'
    merged_df.to_excel(merged_output_file_name, index=False)
    print(f"Merged data saved to {merged_output_file_name}")

    correlation_results = explore_relationships(merged_df)
    t_test_results = perform_two_sample_ttest(merged_df)

    print("Correlation Results:")
    for column, result in correlation_results.items():
        print(f"Column: {column}, Correlation: {result['correlation']:.2f}, P-value: {result['p_value']:.2e}")

    create_presentation(desc_stats, outliers, correlation_results, t_test_results)

"""
    print("Descriptive Statistics:")
    print(desc_stats)
    print("Outliers:")
    print(outliers)
"""

"""
suggestion to save missing values: fill out the missing gaps with classifiers or median values
"""
if __name__ == "__main__":
    main()









